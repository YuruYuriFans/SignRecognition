"""
COMP3419 Assignment Presentation Generator
Generates a PowerPoint presentation based on the traffic sign recognition report

Requirements:
pip install python-pptx

Usage:
python generate_presentation.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """Generate the complete presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Define color scheme
    DARK_BLUE = RGBColor(0, 51, 102)
    LIGHT_BLUE = RGBColor(51, 153, 255)
    GREEN = RGBColor(0, 176, 80)
    RED = RGBColor(192, 0, 0)
    GRAY = RGBColor(89, 89, 89)
    
    # Slide 1: Title Slide
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
    title_frame = title.text_frame
    title_frame.text = "Lightweight CNNs for Traffic Sign Recognition"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = DARK_BLUE
    title_para.alignment = PP_ALIGN.CENTER
    
    subtitle = slide.shapes.add_textbox(Inches(1), Inches(3.2), Inches(8), Inches(0.6))
    subtitle_frame = subtitle.text_frame
    subtitle_frame.text = "A Comparative Study of LeNet-5 vs Modern Architectures"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = LIGHT_BLUE
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    info = slide.shapes.add_textbox(Inches(1), Inches(5), Inches(8), Inches(1.5))
    info_frame = info.text_frame
    info_frame.text = "COMP3419 - Graphics and Multimedia\nBoyan CAI\nNovember 2025"
    for para in info_frame.paragraphs:
        para.font.size = Pt(20)
        para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Problem & Motivation
    slide = add_title_slide(prs, "Problem & Motivation", DARK_BLUE)
    
    left, top = Inches(0.5), Inches(1.8)
    width, height = Inches(9), Inches(5)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    add_bullet(tf, "Challenge: High computational requirements limit embedded traffic sign recognition", 0, 24)
    add_bullet(tf, "Gap: Speed is under-researched despite being crucial for deployment", 0, 24)
    add_bullet(tf, "Research Question: Can the classic LeNet-5 compete with modern lightweight models?", 0, 24, LIGHT_BLUE, True)
    add_bullet(tf, "Focus: GTSRB benchmark with 4 CNN architectures", 0, 24)
    
    # Slide 3: Research Approach
    slide = add_title_slide(prs, "Research Approach", DARK_BLUE)
    
    # Three columns
    col_width = Inches(2.8)
    col_height = Inches(4.5)
    top = Inches(1.8)
    
    # Column 1
    box1 = slide.shapes.add_textbox(Inches(0.5), top, col_width, col_height)
    tf1 = box1.text_frame
    add_subtitle(tf1, "Model Comparison", 22, DARK_BLUE, True)
    add_bullet(tf1, "LeNet-5", 1, 18)
    add_bullet(tf1, "MiniVGG", 1, 18)
    add_bullet(tf1, "MobileNetV2 (α=0.25)", 1, 18)
    add_bullet(tf1, "MobileNetV4-Small", 1, 18)
    
    # Column 2
    box2 = slide.shapes.add_textbox(Inches(3.6), top, col_width, col_height)
    tf2 = box2.text_frame
    add_subtitle(tf2, "Fine-tuning", 22, DARK_BLUE, True)
    add_bullet(tf2, "Label smoothing", 1, 18)
    add_bullet(tf2, "Adam optimizer", 1, 18)
    add_bullet(tf2, "Batch size optimization", 1, 18)
    add_bullet(tf2, "Learning rate tuning", 1, 18)
    
    # Column 3
    box3 = slide.shapes.add_textbox(Inches(6.7), top, col_width, col_height)
    tf3 = box3.text_frame
    add_subtitle(tf3, "Ablation Studies", 22, DARK_BLUE, True)
    add_bullet(tf3, "Data augmentation", 1, 18)
    add_bullet(tf3, "Kernel sizes", 1, 18)
    add_bullet(tf3, "Network depth", 1, 18)
    add_bullet(tf3, "Activation functions", 1, 18)
    
    # Slide 4: Dataset & Metrics
    slide = add_title_slide(prs, "Dataset & Evaluation Metrics", DARK_BLUE)
    
    # Left side - Dataset
    box_left = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4), Inches(4.5))
    tf_left = box_left.text_frame
    add_subtitle(tf_left, "GTSRB Dataset", 24, DARK_BLUE, True)
    add_bullet(tf_left, "50,000+ daytime images", 1, 20)
    add_bullet(tf_left, "43 German traffic sign classes", 1, 20)
    add_bullet(tf_left, "Varying size, lighting, perspective", 1, 20)
    add_bullet(tf_left, "Standard benchmark", 1, 20)
    
    # Right side - Metrics
    box_right = slide.shapes.add_textbox(Inches(5), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_right = box_right.text_frame
    add_subtitle(tf_right, "Evaluation Metrics", 24, DARK_BLUE, True)
    add_bullet(tf_right, "Accuracy: Primary measure", 1, 18)
    add_bullet(tf_right, "Weighted F1/Prec/Rec: Handle imbalance", 1, 18)
    add_bullet(tf_right, "Log Loss: Penalize overconfidence", 1, 18)
    add_bullet(tf_right, "FLOPs: Theoretical complexity", 1, 18)
    add_bullet(tf_right, "Throughput: Real-world speed", 1, 18)
    
    # Slide 5: Main Results
    slide = add_title_slide(prs, "Main Results - Performance Comparison", DARK_BLUE)
    
    # Table data
    table_data = [
        ["Model", "Test Acc", "LogLoss", "FLOPs (G)", "Throughput\n(img/s)", "W-F1"],
        ["MiniVGG", "98.08%", "0.088", "0.060", "29,718", "0.9805"],
        ["LeNet", "96.27%", "0.200", "0.007", "120,504", "0.9619"],
        ["MobileNetV4", "94.44%", "0.234", "0.011", "18,864", "0.9442"],
        ["MobileNetV2", "91.35%", "0.387", "0.002", "52,964", "0.9134"]
    ]
    
    # Add table
    rows, cols = len(table_data), len(table_data[0])
    left, top = Inches(0.8), Inches(2)
    width, height = Inches(8.4), Inches(2.5)
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    # Fill table
    for i, row_data in enumerate(table_data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_data
            
            # Format header
            if i == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = DARK_BLUE
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(16)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    para.alignment = PP_ALIGN.CENTER
            else:
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(16)
                    para.alignment = PP_ALIGN.CENTER
                    
                # Highlight best values
                if (j == 1 and i == 1) or (j == 2 and i == 1):  # MiniVGG accuracy & log loss
                    para.font.color.rgb = GREEN
                    para.font.bold = True
                elif (j == 3 and i == 2) or (j == 4 and i == 2):  # LeNet FLOPs & throughput
                    para.font.color.rgb = GREEN
                    para.font.bold = True
    
    # Key insights
    insights_box = slide.shapes.add_textbox(Inches(0.8), Inches(4.8), Inches(8.4), Inches(1.5))
    tf_insights = insights_box.text_frame
    add_subtitle(tf_insights, "Key Insights:", 20, DARK_BLUE, True)
    add_bullet(tf_insights, "MiniVGG: Highest accuracy but 8.6× more FLOPs", 1, 18)
    add_bullet(tf_insights, "LeNet-5: 4× faster than MiniVGG with acceptable accuracy", 1, 18)
    add_bullet(tf_insights, "Clear accuracy-speed trade-off demonstrated", 1, 18)
    
    # Slide 6: Ablation Studies
    slide = add_title_slide(prs, "Ablation Study Results", DARK_BLUE)
    
    # Three mini-tables side by side
    table_width = Inches(2.8)
    table_height = Inches(2.2)
    top = Inches(2)
    
    # Data Augmentation Table
    aug_data = [
        ["Aug Type", "LeNet Acc"],
        ["None", "99.64%"],
        ["Basic", "99.82%"],
        ["Advanced", "99.77%"]
    ]
    add_mini_table(slide, aug_data, Inches(0.5), top, table_width, table_height, 
                   "Data Augmentation", "No significant difference", DARK_BLUE)
    
    # Kernel Sizes Table
    kernel_data = [
        ["Config", "Acc"],
        ["5×5 (Base)", "96.27%"],
        ["All 7×7", "39.90%"],
        ["All 3×3", "29.93%"]
    ]
    add_mini_table(slide, kernel_data, Inches(3.6), top, table_width, table_height,
                   "Kernel Sizes", "Original optimal", DARK_BLUE, GREEN)
    
    # Activation Functions Table
    act_data = [
        ["Function", "Acc"],
        ["ReLU", "96.27%"],
        ["LeakyReLU", "36.83%"],
        ["GELU", "33.43%"]
    ]
    add_mini_table(slide, act_data, Inches(6.7), top, table_width, table_height,
                   "Activation Functions", "ReLU optimal", DARK_BLUE, GREEN)
    
    # Key findings
    findings_box = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.2))
    tf_findings = findings_box.text_frame
    add_subtitle(tf_findings, "Key Findings:", 20, DARK_BLUE, True)
    add_bullet(tf_findings, "LeNet-5's original design is optimal for shallow architecture", 1, 18)
    add_bullet(tf_findings, "GTSRB has sufficient natural variation (augmentation minimal impact)", 1, 18)
    
    # Slide 7: Fine-tuning Results
    slide = add_title_slide(prs, "Fine-tuning Results", DARK_BLUE)
    
    content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4.5))
    tf_content = content_box.text_frame
    
    add_subtitle(tf_content, "Techniques Tested:", 24, DARK_BLUE, True)
    add_bullet(tf_content, "Label smoothing regularization", 1, 20)
    add_bullet(tf_content, "Adam optimizer", 1, 20)
    add_bullet(tf_content, "Various batch sizes (32, 64, 128)", 1, 20)
    add_bullet(tf_content, "Learning rate adjustments", 1, 20)
    
    tf_content.add_paragraph()
    add_subtitle(tf_content, "Result: No statistically significant improvement", 24, RED, True)
    
    tf_content.add_paragraph()
    add_subtitle(tf_content, "Why?", 22, DARK_BLUE, True)
    add_bullet(tf_content, "LeNet-5's lightweight architecture operates near its performance ceiling", 1, 20, LIGHT_BLUE)
    add_bullet(tf_content, "Validates robustness of original architecture", 1, 20, LIGHT_BLUE)
    
    # Slide 8: Key Insights
    slide = add_title_slide(prs, "Key Insights: Why These Results?", DARK_BLUE)
    
    # Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_left = left_box.text_frame
    add_subtitle(tf_left, "Why MiniVGG is More Accurate", 20, DARK_BLUE, True)
    add_bullet(tf_left, "Deeper architecture (more layers)", 1, 16)
    add_bullet(tf_left, "Sequential 3×3 conv blocks", 1, 16)
    add_bullet(tf_left, "Robust feature hierarchy", 1, 16)
    add_bullet(tf_left, "Captures complex patterns", 1, 16)
    tf_left.add_paragraph()
    add_bullet(tf_left, "Cost: 8.6× more computation", 1, 18, RED, True)
    
    # Right column
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_right = right_box.text_frame
    add_subtitle(tf_right, "Why LeNet-5 is Faster", 20, DARK_BLUE, True)
    add_bullet(tf_right, "Shallow 2-layer design", 1, 16)
    add_bullet(tf_right, "Efficient 5×5 kernels", 1, 16)
    add_bullet(tf_right, "Minimal computational overhead", 1, 16)
    add_bullet(tf_right, "Adequate representational capacity", 1, 16)
    tf_right.add_paragraph()
    add_bullet(tf_right, "Benefit: 4× higher throughput", 1, 18, GREEN, True)
    
    # Bottom banner
    banner = slide.shapes.add_textbox(Inches(1), Inches(6.5), Inches(8), Inches(0.8))
    tf_banner = banner.text_frame
    banner_para = tf_banner.paragraphs[0]
    banner_para.text = "Architecture depth directly impacts accuracy-speed trade-off"
    banner_para.font.size = Pt(22)
    banner_para.font.bold = True
    banner_para.font.color.rgb = LIGHT_BLUE
    banner_para.alignment = PP_ALIGN.CENTER
    
    # Slide 9: Practical Implications
    slide = add_title_slide(prs, "Practical Implications: When to Use Each Model?", DARK_BLUE)
    
    # Left: LeNet-5
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(2), Inches(4.5), Inches(4))
    tf_left = left_box.text_frame
    add_subtitle(tf_left, "Choose LeNet-5 when:", 22, GREEN, True)
    add_bullet(tf_left, "Real-time processing is critical", 1, 18)
    add_bullet(tf_left, "Hardware resources are limited", 1, 18)
    add_bullet(tf_left, "96% accuracy is acceptable", 1, 18)
    tf_left.add_paragraph()
    add_subtitle(tf_left, "Use Case:", 20, DARK_BLUE, True)
    add_bullet(tf_left, "Embedded dashcams, edge devices", 1, 18, GRAY)
    
    # Right: MiniVGG
    right_box = slide.shapes.add_textbox(Inches(5), Inches(2), Inches(4.5), Inches(4))
    tf_right = right_box.text_frame
    add_subtitle(tf_right, "Choose MiniVGG when:", 22, GREEN, True)
    add_bullet(tf_right, "Maximum accuracy is required", 1, 18)
    add_bullet(tf_right, "Computational resources available", 1, 18)
    add_bullet(tf_right, "Safety-critical applications", 1, 18)
    tf_right.add_paragraph()
    add_subtitle(tf_right, "Use Case:", 20, DARK_BLUE, True)
    add_bullet(tf_right, "Autonomous vehicles, servers", 1, 18, GRAY)
    
    # Slide 10: Limitations & Future Work
    slide = add_title_slide(prs, "Limitations & Future Work", DARK_BLUE)
    
    # Left: Limitations
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_left = left_box.text_frame
    add_subtitle(tf_left, "Current Limitations", 22, RED, True)
    add_bullet(tf_left, "Daytime images only", 1, 18)
    add_bullet(tf_left, "Single-label classification", 1, 18)
    add_bullet(tf_left, "Ideal conditions only", 1, 18)
    add_bullet(tf_left, "Limited labeled data scenarios", 1, 18)
    
    # Right: Future Directions
    right_box = slide.shapes.add_textbox(Inches(5), Inches(1.8), Inches(4.5), Inches(4.5))
    tf_right = right_box.text_frame
    add_subtitle(tf_right, "Future Directions", 22, GREEN, True)
    add_bullet(tf_right, "Test on adverse conditions (night, weather)", 1, 18)
    add_bullet(tf_right, "Multi-label classification (SRN)", 1, 18)
    add_bullet(tf_right, "GAN-based data synthesis", 1, 18)
    add_bullet(tf_right, "Real-world deployment testing", 1, 18)
    add_bullet(tf_right, "Cross-dataset validation", 1, 18)
    
    # Slide 11: Conclusions
    slide = add_title_slide(prs, "Conclusions", DARK_BLUE)
    
    content_box = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(4.5))
    tf = content_box.text_frame
    
    add_subtitle(tf, "Key Findings:", 24, DARK_BLUE, True)
    
    # LeNet-5 Performance
    p1 = tf.add_paragraph()
    p1.text = "✓ LeNet-5: 96.27% accuracy with 120,504 img/s throughput"
    p1.level = 1
    p1.font.size = Pt(20)
    
    p2 = tf.add_paragraph()
    p2.text = "✓ Remains viable for embedded systems after 30 years"
    p2.level = 1
    p2.font.size = Pt(20)
    
    p3 = tf.add_paragraph()
    p3.text = "✓ Original design choices validated through ablation"
    p3.level = 1
    p3.font.size = Pt(20)
    
    tf.add_paragraph()
    
    # Main insight
    insight_para = tf.add_paragraph()
    insight_para.text = "Main Insight: Fundamental accuracy-speed trade-off"
    insight_para.font.size = Pt(22)
    insight_para.font.bold = True
    insight_para.font.color.rgb = LIGHT_BLUE
    
    detail_para = tf.add_paragraph()
    detail_para.text = "MiniVGG achieves 98% accuracy but LeNet-5 delivers 4× faster inference with 87% lower computational cost"
    detail_para.level = 1
    detail_para.font.size = Pt(18)
    
    tf.add_paragraph()
    
    # Recommendation
    rec_para = tf.add_paragraph()
    rec_para.text = "Recommendation:"
    rec_para.font.size = Pt(22)
    rec_para.font.bold = True
    rec_para.font.color.rgb = DARK_BLUE
    
    final_para = tf.add_paragraph()
    final_para.text = "LeNet-5 is optimal for resource-constrained real-time traffic sign recognition"
    final_para.level = 1
    final_para.font.size = Pt(20)
    final_para.font.color.rgb = GREEN
    final_para.font.bold = True
    
    # Slide 12: Thank You
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    thanks = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(6), Inches(1))
    thanks_frame = thanks.text_frame
    thanks_para = thanks_frame.paragraphs[0]
    thanks_para.text = "Thank You"
    thanks_para.font.size = Pt(54)
    thanks_para.font.bold = True
    thanks_para.font.color.rgb = DARK_BLUE
    thanks_para.alignment = PP_ALIGN.CENTER
    
    contact = slide.shapes.add_textbox(Inches(2), Inches(4), Inches(6), Inches(1.5))
    contact_frame = contact.text_frame
    contact_frame.text = "Questions?\n\nCode & Results: [GitHub Repository]"
    for para in contact_frame.paragraphs:
        para.font.size = Pt(24)
        para.alignment = PP_ALIGN.CENTER
    
    return prs

def add_title_slide(prs, title_text, color):
    """Add a slide with title"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title_text
    title_para.font.size = Pt(36)
    title_para.font.bold = True
    title_para.font.color.rgb = color
    
    # Add line under title
    line = slide.shapes.add_shape(1, Inches(0.5), Inches(1.4), Inches(9), Inches(0))
    line.line.color.rgb = color
    line.line.width = Pt(2)
    
    return slide

def add_bullet(text_frame, text, level, size, color=None, bold=False):
    """Add a bullet point to text frame"""
    p = text_frame.add_paragraph()
    p.text = text
    p.level = level
    p.font.size = Pt(size)
    p.font.bold = bold
    if color:
        p.font.color.rgb = color

def add_subtitle(text_frame, text, size, color, bold=True):
    """Add a subtitle"""
    p = text_frame.paragraphs[0] if len(text_frame.paragraphs) > 0 else text_frame.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    text_frame.add_paragraph()

def add_mini_table(slide, data, left, top, width, height, title, note, title_color, highlight_color=None):
    """Add a small table with title and note"""
    # Title
    title_box = slide.shapes.add_textbox(left, top - Inches(0.4), width, Inches(0.3))
    tf_title = title_box.text_frame
    p_title = tf_title.paragraphs[0]
    p_title.text = title
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = title_color
    p_title.alignment = PP_ALIGN.CENTER
    
    # Table
    rows, cols = len(data), len(data[0])
    table = slide.shapes.add_table(rows, cols, left, top, width, height).table
    
    for i, row_data in enumerate(data):
        for j, cell_data in enumerate(row_data):
            cell = table.cell(i, j)
            cell.text = cell_data
            
            if i == 0:  # Header
                cell.fill.solid()
                cell.fill.fore_color.rgb = title_color
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(14)
                    para.font.bold = True
                    para.font.color.rgb = RGBColor(255, 255, 255)
                    para.alignment = PP_ALIGN.CENTER
            else:
                for para in cell.text_frame.paragraphs:
                    para.font.size = Pt(14)
                    para.alignment = PP_ALIGN.CENTER
                    
                # Highlight first data row if color specified
                if i == 1 and highlight_color:
                    para.font.color.rgb = highlight_color
                    para.font.bold = True
    
    # Note
    note_box = slide.shapes.add_textbox(left, top + height + Inches(0.05), width, Inches(0.3))
    tf_note = note_box.text_frame
    p_note = tf_note.paragraphs[0]
    p_note.text = note
    p_note.font.size = Pt(12)
    p_note.font.italic = True
    p_note.alignment = PP_ALIGN.CENTER

if __name__ == "__main__":
    print("Generating COMP3419 presentation...")
    prs = create_presentation()
    
    output_file = "COMP3419_Traffic_Sign_Recognition_Presentation.pptx"
    prs.save(output_file)
    print(f"Presentation saved as: {output_file}")
    print("\nSlide Summary:")
    print("1. Title Slide")
    print("2. Problem & Motivation")
    print("3. Research Approach")
    print("4. Dataset & Evaluation Metrics")
    print("5. Main Results - Performance Comparison")
    print("6. Ablation Study Results")
    print("7. Fine-tuning Results")
    print("8. Key Insights: Why These Results?")
    print("9. Practical Implications")
    print("10. Limitations & Future Work")
    print("11. Conclusions")
    print("12. Thank You")
    print("\nTotal: 12 slides")
    print("Estimated presentation time: 4-5 minutes")